from doctomarkdown.base import BaseConverter, PageResult
from pptx import Presentation
import tempfile
import os
import shutil

class PptxToMarkdown(BaseConverter):
    """Converter for PPTX files to Markdown format with optional LLM post-processing."""

    def __init__(
        self, filepath, extract_images=False, extract_tables=False,
        output_path=None, llm_client=None, llm_model=None,
        system_prompt=None, user_prompt_template=None,
        output_type='markdown', **kwargs
    ):
        super().__init__(
            filepath=filepath,
            extract_images=extract_images,
            extract_tables=extract_tables,
            output_path=output_path,
            llm_client=llm_client,
            llm_model=llm_model,
            output_type=output_type,
            **kwargs
        )
        self.system_prompt = system_prompt or "You are a helpful assistant that converts slide content into Markdown."
        self.user_prompt_template = user_prompt_template or "Convert the following presentation content into Markdown:\n\n{content}"

    def call_llm(self, content):
        """Call the LLM with the configured prompt template."""
        prompt = self.user_prompt_template.format(content=content)
        response = self.llm_client.chat.completions.create(
            model=self.llm_model,
            messages=[
                {"role": "system", "content": self.system_prompt},
                {"role": "user", "content": prompt}
            ]
        )
        return response.choices[0].message.content

    def extract_content(self):
        temp_dir = None

        # Try PDF fallback if LLM client is provided
        try:
            print(f"[INFO] Starting PPTX extraction for: {self.filepath}")
            if self.llm_client:
                from doctomarkdown.utils.pptx_to_pdf import convert_pptx_to_pdf
                from doctomarkdown.converters.pdf_to_markdown import PdfToMarkdown

                temp_dir = tempfile.mkdtemp()
                input_dir = os.path.dirname(self.filepath)
                output_dir = temp_dir

                convert_pptx_to_pdf(input_dir, output_dir)
                pptx_basename = os.path.splitext(os.path.basename(self.filepath))[0]
                pdf_path = os.path.join(output_dir, pptx_basename + '.pdf')

                if os.path.exists(pdf_path):
                    pdf_converter = PdfToMarkdown(
                        filepath=pdf_path,
                        llm_client=self.llm_client,
                        llm_model=self.llm_model,
                        system_prompt=self.system_prompt,
                        user_prompt_template=self.user_prompt_template,
                        extract_images=self.extract_images,
                        extract_tables=self.extract_tables,
                        output_path=None,
                        output_type=self.output_type
                    )
                    pages = pdf_converter.extract_content()
                    self._markdown = getattr(pdf_converter, '_markdown', None)
                    print(f"[SUCCESS] Extraction is successful via PDF fallback for: {self.filepath}")
                    return pages
                else:
                    print(f"[ERROR] PDF file not found after conversion: {pdf_path}")
        except Exception as e:
            print(f"[FAILURE] Error during PPTX to PDF conversion: {e}")
        finally:
            if temp_dir and os.path.exists(temp_dir):
                shutil.rmtree(temp_dir)

        # Fallback to standard slide-by-slide processing
        try:
            doc = Presentation(self.filepath)
            pages = []
            markdown_lines = []

            for page_number, slide in enumerate(doc.slides, 1):
                slide_content = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text.strip())

                page_text = "\n".join(slide_content)
                pages.append(PageResult(page_number, page_text))
                markdown_lines.append(f"Page Number: {page_number}\nPage Content:\n{page_text}\n")

            full_markdown = "\n".join(markdown_lines)

            # LLM post-processing on full presentation
            if self.llm_client and self.llm_model and full_markdown.strip():
                try:
                    full_markdown = self.call_llm(full_markdown)
                except Exception as e:
                    print(f"[WARNING] LLM post-processing failed: {e}")

            self._markdown = full_markdown
            print(f"[SUCCESS] Extraction succeeded using standard PPTX logic for: {self.filepath}")
            return pages

        except Exception as e:
            print(f"[FAILURE] Error during standard PPTX extraction: {e}")
            raise
